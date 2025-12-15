from copy import deepcopy
from pathlib import Path
from io import BytesIO
import random

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches
from pptx.oxml.xmlchemy import OxmlElement


def duplicate_slide(prs: Presentation, source_slide):
    """Create a copy of the slide at `index` and append it to the deck."""
    source = source_slide
    layout = source.slide_layout
    
    # Create new slide using layout
    new_slide = prs.slides.add_slide(layout)

    # Map old rel IDs to new ones so embedded images/hyperlinks work.
    rel_id_map = {}
    for rel_obj in list(source.part.rels._rels.values()):
        if "notesSlide" in rel_obj.reltype:
            continue
        new_rid = new_slide.part.relate_to(rel_obj._target, rel_obj.reltype)
        rel_id_map[rel_obj.rId] = new_rid

    # Copy background from source slide
    source_csld = source.element.find(qn("p:cSld"))
    new_csld = new_slide.element.find(qn("p:cSld"))
    
    if source_csld is not None and new_csld is not None:
        source_bg = source_csld.find(qn("p:bg"))
        if source_bg is not None:
            # Remove any existing background in new slide
            existing_bg = new_csld.find(qn("p:bg"))
            if existing_bg is not None:
                new_csld.remove(existing_bg)
            
            # Deep copy the background element
            new_bg = deepcopy(source_bg)
            
            # Insert background as first element in cSld (before spTree)
            # Get all children to find insertion point
            children = list(new_csld)
            if children:
                # Insert before first child
                new_csld.insert(0, new_bg)
            else:
                new_csld.append(new_bg)

    # Copy shapes from source
    for shape in source.shapes:
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    # Normalize logo position: move logo to top-right corner
    logo_top = Inches(0)
    slide_w = prs.slide_width
    margin_right = Inches(0.3)  # Small margin from right edge

    for shp in new_slide.shapes:
        try:
            # Check for small pictures (logos)
            if shp.shape_type == 13 and shp.width <= Inches(3) and shp.height <= Inches(3):
                # Position logo at top-right
                shp.top = logo_top
                shp.left = slide_w - shp.width - margin_right
            # Check for FREEFORM shapes with embedded images in top-right area (logos)
            elif shp.shape_type == 5:  # FREEFORM
                # Check if it has an embedded image and is positioned near top-right
                has_image = False
                for blip in shp.element.iter(qn("a:blip")):
                    if blip.get(qn("r:embed")):
                        has_image = True
                        break
                if has_image and shp.top < Inches(2) and shp.left > slide_w * 0.7:
                    # This is likely the logo - position at top-right
                    shp.top = logo_top
                    shp.left = slide_w - shp.width - margin_right
        except Exception:
            continue

    # Update relationship references inside the new slide XML
    rel_key = qn("r:embed")
    hyperlink_key = qn("r:id")
    for blip in new_slide.element.iter(qn("a:blip")):
        rid = blip.get(rel_key)
        if rid and rid in rel_id_map:
            blip.set(rel_key, rel_id_map[rid])
    for elm in new_slide.element.iter():
        rid = elm.get(hyperlink_key)
        if rid and rid in rel_id_map:
            elm.set(hyperlink_key, rel_id_map[rid])

    # Copy notes if present
    if source.has_notes_slide:
        new_slide.notes_slide.notes_text_frame.text = source.notes_slide.notes_text_frame.text

    return new_slide


def copy_slide_from_other(prs: Presentation, source_slide):
    """Copy a slide from another presentation into prs, preserving layout name, background,
    shapes, and pictures. Returns the new slide."""
    # Find matching layout by name
    layout = None
    for lo in prs.slide_layouts:
        if lo.name == source_slide.slide_layout.name:
            layout = lo
            break
    if layout is None:
        layout = prs.slide_layouts[0]

    new_slide = prs.slides.add_slide(layout)

    # Copy background
    source_csld = source_slide.element.find(qn("p:cSld"))
    new_csld = new_slide.element.find(qn("p:cSld"))
    if source_csld is not None and new_csld is not None:
        source_bg = source_csld.find(qn("p:bg"))
        if source_bg is not None:
            existing_bg = new_csld.find(qn("p:bg"))
            if existing_bg is not None:
                new_csld.remove(existing_bg)
            new_bg = deepcopy(source_bg)
            children = list(new_csld)
            if children:
                new_csld.insert(0, new_bg)
            else:
                new_csld.append(new_bg)

    # Copy shapes; for pictures, re-add with blob to ensure media is present
    for shape in source_slide.shapes:
        try:
            if shape.shape_type == 13:  # PICTURE
                img_stream = BytesIO(shape.image.blob)
                pic = new_slide.shapes.add_picture(img_stream, shape.left, shape.top, width=shape.width, height=shape.height)
                # Preserve alt text/title if any
                pic.alt_text = shape.alt_text
            else:
                new_el = deepcopy(shape.element)
                new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
        except Exception:
            continue

    return new_slide


def remove_slide(prs: Presentation, index: int):
    """Remove slide at index and its relationships cleanly to avoid orphan parts."""
    slide = prs.slides[index]
    slide_id = prs.slides._sldIdLst[index]
    r_id = slide_id.rId
    prs.part.drop_rel(r_id)
    prs.slides._sldIdLst.remove(slide_id)


def _paragraph_level(p):
    """Extract outline/numbering level from a docx paragraph if present."""
    try:
        pPr = p._p.pPr
        if pPr is None:
            return 0
        
        # First check for explicit numbering level
        numPr = pPr.numPr
        if numPr is not None:
            ilvl = numPr.ilvl
            if ilvl is not None and ilvl.val is not None:
                return int(ilvl.val)
        
        # If no numbering, check for indentation (left indent)
        # Indentation is often used to indicate hierarchy
        ind = pPr.ind
        if ind is not None and ind.left is not None:
            # Convert indentation to level (every 720 twips = 0.5 inches = 1 level)
            # Typical indentation: 0 = level 0, 720+ = level 1, 1440+ = level 2, etc.
            left_indent = int(ind.left) if hasattr(ind.left, '__int__') else 0
            if left_indent > 0:
                # Estimate level based on indentation (720 twips per level)
                level = min(int(left_indent / 720), 2)  # Cap at level 2
                return level
    except Exception:
        return 0
    return 0


def parse_docx(doc_path: Path):
    """Parse Mod1.docx into a list of slides with title and bullet items including levels."""
    from docx import Document

    doc = Document(doc_path)
    slides = []
    current = None
    for p in doc.paragraphs:
        text = p.text.strip()
        if not text:
            continue
        if text.startswith("Slide "):
            # start new slide
            if current:
                slides.append(current)
            parts = text.split("–", 1)
            title = parts[1].strip() if len(parts) > 1 else text
            current = {"title": title, "bullets": []}
        else:
            if current is None:
                continue
            level = _paragraph_level(p)
            current["bullets"].append({"text": text, "level": level})
    if current:
        slides.append(current)
    return slides


def apply_content(new_slide, title_text: str, bullet_items):
    """Apply title and bullet text to duplicated slide while keeping formatting."""
    def ensure_bullet(para):
        """Force bullet display on a paragraph using the slide's default bullet style."""
        try:
            para.level = 0
        except Exception:
            pass
        pPr = para._p.get_or_add_pPr()
        # Remove buNone if present so the template's bullet style applies.
        for bu_none in list(pPr.findall(qn("a:buNone"))):
            pPr.remove(bu_none)

    # Find title and body shapes based on slide 2 structure, with fallbacks
    title_shape = None
    body_shape = None
    text_shapes = []
    for shape in new_slide.shapes:
        if hasattr(shape, "text"):
            text_shapes.append(shape)
            if "OBJECTIVES" in shape.text.upper():
                title_shape = shape
            elif "By the end of this module" in shape.text:
                body_shape = shape

    # Fallback: if body_shape not found, pick the longest text box (excluding title)
    if body_shape is None:
        longest = None
        longest_len = -1
        for shp in text_shapes:
            if shp is title_shape:
                continue
            txt = getattr(shp, "text", "") or ""
            if len(txt) > longest_len:
                longest = shp
                longest_len = len(txt)
        body_shape = longest

    # Update title while preserving formatting by reusing existing run
    if title_shape is not None and hasattr(title_shape, "text_frame"):
        # Expand width to minimize wrapping
        try:
            title_shape.width = Inches(15.54)
        except Exception:
            pass
        p = title_shape.text_frame.paragraphs[0]
        runs = list(p.runs)
        if runs:
            for extra in runs[1:]:
                p._p.remove(extra._r)
            # set size to ~42pt to help fit in single line
            from pptx.util import Pt
            runs[0].font.size = Pt(42)
            runs[0].text = title_text
        else:
            p.text = title_text

    # Update body text preserving bullet formatting (reuse existing paragraphs/runs)
    if body_shape is not None and hasattr(body_shape, "text_frame"):
        tf = body_shape.text_frame
        paras = list(tf.paragraphs)
        template_para = paras[-1] if paras else None

        # Ensure we have enough paragraphs by cloning the last one
        while len(tf.paragraphs) < len(bullet_items) and template_para is not None:
            new_p_el = deepcopy(template_para._p)
            tf._element.append(new_p_el)

        # Update text in place to keep formatting
        for idx, item in enumerate(bullet_items):
            if idx >= len(tf.paragraphs):
                break
            text = item["text"]
            lvl = item.get("level", 0)
            para = tf.paragraphs[idx]
            runs = list(para.runs)
            if runs:
                for extra in runs[1:]:
                    para._p.remove(extra._r)
                runs[0].text = text
            else:
                para.text = text
            
            # Set paragraph level - only apply indentation if level is explicitly > 0
            # Don't add indentation for level 0 (normal level)
            try:
                if lvl > 0:
                    para.level = lvl
                    # Set explicit indentation for sub-levels
                    pPr = para._p.get_or_add_pPr()
                    # Set left indent: 720000 EMUs per level (0.5 inches)
                    indent_emus = lvl * 720000
                    ind = pPr.find(qn("a:ind"))
                    if ind is None:
                        from pptx.oxml import parse_xml
                        ind_str = f'<a:ind xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" l="{indent_emus}"/>'
                        ind_el = parse_xml(ind_str)
                        pPr.append(ind_el)
                    else:
                        ind.set("l", str(indent_emus))
                else:
                    # For level 0, ensure no indentation is set
                    para.level = 0
                    # Remove any existing indentation
                    try:
                        pPr = para._p.get_or_add_pPr()
                        ind = pPr.find(qn("a:ind"))
                        if ind is not None:
                            pPr.remove(ind)
                    except:
                        pass
            except Exception:
                # If level setting fails, just ensure level 0 has no indentation
                try:
                    if lvl == 0:
                        pPr = para._p.get_or_add_pPr()
                        ind = pPr.find(qn("a:ind"))
                        if ind is not None:
                            pPr.remove(ind)
                except:
                    pass
            ensure_bullet(para)

        # Remove extra paragraphs beyond bullet_lines
        while len(tf.paragraphs) > len(bullet_items):
            tf._element.remove(tf.paragraphs[-1]._p)

    # Position logo at top-right corner
    # Get slide width from the presentation
    try:
        prs = new_slide.part.package.presentation_part.presentation
        slide_w = prs.slide_width
    except:
        slide_w = Inches(10)  # Default fallback
    
    margin_right = Inches(0.3)
    logo_top = Inches(0)
    
    for shp in new_slide.shapes:
        try:
            # Position small pictures (logos) at top-right
            if shp.shape_type == 13 and shp.width <= Inches(3) and shp.height <= Inches(3):
                shp.top = logo_top
                shp.left = slide_w - shp.width - margin_right
            # Check for FREEFORM shapes with embedded images in top-right area (logos)
            elif shp.shape_type == 5:  # FREEFORM
                # Check if it has an embedded image and is positioned near top-right
                has_image = False
                for blip in shp.element.iter(qn("a:blip")):
                    if blip.get(qn("r:embed")):
                        has_image = True
                        break
                if has_image and shp.top < Inches(2) and shp.left > slide_w * 0.7:
                    # This is likely the logo - position at top-right
                    shp.top = logo_top
                    shp.left = slide_w - shp.width - margin_right
        except Exception:
            continue
    
    # Position decorative bars at a fixed top offset (1.06")
    # But exclude logos (FREEFORM shapes with images in top-right area)
    for shp in new_slide.shapes:
        try:
            if shp.shape_type == 5 and shp.top < Inches(4):  # FREEFORM near top
                # Check if this is a logo (has embedded image and is in top-right)
                is_logo = False
                try:
                    has_image = False
                    for blip in shp.element.iter(qn("a:blip")):
                        if blip.get(qn("r:embed")):
                            has_image = True
                            break
                    if has_image and shp.left > slide_w * 0.7:
                        is_logo = True
                except:
                    pass
                if not is_logo:
                    shp.top = Inches(1.06)
        except Exception:
            continue


def main():
    root = Path(__file__).resolve().parent
    src_path = root / "sample.pptx"
    dst_path = root / "sample_generated.pptx"
    doc_path = root / "Mod1.docx"
    elements_dir = root / "elements"

    slides_data = parse_docx(doc_path)
    prs = Presentation(src_path)

    if len(prs.slides) < 2:
        raise ValueError("Presentation needs at least 2 slides to use as template.")

    template_index = 1  # original slide 2
    template_slide = prs.slides[template_index]
    slide3_index = 2 if len(prs.slides) > 2 else None

    generated_slides = []
    for slide_info in slides_data:
        target_slide = duplicate_slide(prs, template_slide)
        apply_content(target_slide, slide_info["title"], slide_info["bullets"])
        generated_slides.append(target_slide)

    # Load shapes from shapes folder (for decorative elements)
    shapes_dir = root / "shapes"
    slide3_shapes = []
    if shapes_dir.exists() and shapes_dir.is_dir():
        for p in sorted(shapes_dir.iterdir()):
            # Skip SVG files as python-pptx doesn't support them directly
            if p.suffix.lower() == ".svg":
                continue
            if p.suffix.lower() in [".png", ".jpg", ".jpeg", ".gif", ".bmp"]:
                try:
                    with open(p, "rb") as f:
                        blob = f.read()
                    from PIL import Image
                    # Get image dimensions
                    try:
                        with Image.open(p) as im:
                            px_w, px_h = im.size
                            dpi = im.info.get("dpi", (96, 96))
                            dpi_x = dpi[0] or 96
                            dpi_y = dpi[1] or 96
                            w = Inches(px_w / dpi_x)
                            h = Inches(px_h / dpi_y)
                    except Exception:
                        # Fallback: use default size if image can't be read
                        w = Inches(2)
                        h = Inches(2)
                    slide3_shapes.append({
                        "kind": "pic",
                        "blob": blob,
                        "width": w,
                        "height": h,
                    })
                except Exception:
                    continue

    # remove the original slide 3 (if present) and the template slide to avoid stale content
    if slide3_index is not None:
        remove_slide(prs, slide3_index)
    remove_slide(prs, template_index)

    # Use shapes from shapes folder as primary source
    elements = slide3_shapes
    
    # Also load decorative elements from the elements folder (PNGs only) and add them
    if elements_dir.exists() and elements_dir.is_dir():
        for p in sorted(elements_dir.iterdir()):
            if p.suffix.lower() == ".png":
                try:
                    with open(p, "rb") as f:
                        blob = f.read()
                    from PIL import Image
                    with Image.open(p) as im:
                        px_w, px_h = im.size
                        dpi = im.info.get("dpi", (96, 96))
                        dpi_x = dpi[0] or 96
                        dpi_y = dpi[1] or 96
                        w = Inches(px_w / dpi_x)
                        h = Inches(px_h / dpi_y)
                    elements.append(
                        {
                            "kind": "pic",
                            "blob": blob,
                            "width": w,
                            "height": h,
                        }
                    )
                except Exception:
                    continue

    # Add a decorative element to bottom-right of each slide (keep original size)
    if elements:
        slide_w = prs.slide_width
        slide_h = prs.slide_height
        margin = Inches(0)  # No margin - bottom 0, right 0
        rel_key = qn("r:embed")

        for idx, slide in enumerate(prs.slides):
            if idx == 0:
                continue  # skip decorative element on slide 1
            if not elements:
                print(f"Warning: No elements available for slide {idx + 1}")
                continue  # Ensure elements list is not empty
            # Randomly select an element from the available images
            elem = random.choice(elements)
            w = elem.get("width")
            h = elem.get("height")
            if w is None or h is None:
                # Use default size if unknown
                w = Inches(2)
                h = Inches(2)
            left = slide_w - w - margin
            top = slide_h - h - margin
            if elem["kind"] == "pic":
                try:
                    # Create a fresh BytesIO object for each slide
                    img_stream = BytesIO(elem["blob"])
                    if w and h:
                        slide.shapes.add_picture(img_stream, left, top, width=w, height=h)
                    else:
                        slide.shapes.add_picture(img_stream, left, top)
                except Exception as e:
                    print(f"Warning: Could not add image to slide {idx + 1}: {e}")
                    # Try to add with default size as fallback
                    try:
                        img_stream = BytesIO(elem["blob"])
                        slide.shapes.add_picture(img_stream, left, top, width=Inches(2), height=Inches(2))
                    except Exception:
                        continue
            else:
                try:
                    new_el = deepcopy(elem["element"])
                    # Update position
                    sppr = new_el.find(qn("p:spPr"))
                    if sppr is not None:
                        xfrm = sppr.find(qn("a:xfrm"))
                        if xfrm is not None:
                            off = xfrm.find(qn("a:off"))
                            if off is not None:
                                off.set("x", str(int(left)))
                                off.set("y", str(int(top)))
                            else:
                                # Create transform if it doesn't exist
                                from pptx.oxml import parse_xml
                                xfrm_str = f'<a:xfrm xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:off x="{int(left)}" y="{int(top)}"/><a:ext cx="{int(w)}" cy="{int(h)}"/></a:xfrm>'
                                xfrm_el = parse_xml(xfrm_str)
                                if sppr.find(qn("a:xfrm")) is None:
                                    sppr.insert(0, xfrm_el)
                    slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
                except Exception:
                    continue

    prs.save(dst_path)
    print(f"Saved updated presentation to {dst_path}")


if __name__ == "__main__":
    main()

